#!/usr/bin/env python3
"""Generate the Contxt pitch deck (docs/contxt-pitch.pptx).

Dark Neo-Kinpaku theme (gold = private, patina = shared). Self-contained — run:
    python3 docs/build_deck.py
Then open the .pptx (PowerPoint / Keynote / Google Slides) and File → Export → PDF.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

INK = RGBColor(0x0E, 0x0E, 0x10)
CHAMPAGNE = RGBColor(0xEC, 0xEA, 0xE4)
TEXT = RGBColor(0xCF, 0xCD, 0xC7)
MUTED = RGBColor(0x9A, 0x98, 0x92)
GOLD = RGBColor(0xE0, 0xB6, 0x4D)
PATINA = RGBColor(0x5F, 0xB6, 0xAD)
FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
LEFT = Inches(0.9)
W = Inches(11.5)


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK


def _mark(slide):
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT, Inches(0.55), Inches(0.26), Inches(0.26))
    sq.fill.background()
    sq.line.color.rgb = GOLD
    sq.line.width = Pt(2)
    ln = slide.shapes.add_connector(2, LEFT + Inches(0.26), Inches(0.55), LEFT, Inches(0.81))
    ln.line.color.rgb = GOLD
    ln.line.width = Pt(2)


def _tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _run(p, text, size, color, bold=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = FONT
    return r


def _footer(slide, n):
    tf = _tb(slide, LEFT, Inches(6.95), W, Inches(0.4))
    _run(tf.paragraphs[0], "CONTXT", 10, MUTED, bold=True)
    tf2 = _tb(slide, Inches(11.8), Inches(6.95), Inches(1.2), Inches(0.4))
    pr = tf2.paragraphs[0]
    pr.alignment = PP_ALIGN.RIGHT
    _run(pr, str(n), 10, MUTED)


def content(n, kicker, title, bullets, accent=GOLD):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    _mark(s)
    kf = _tb(s, LEFT, Inches(1.0), W, Inches(0.4))
    _run(kf.paragraphs[0], kicker.upper(), 13, accent, bold=True)
    ttf = _tb(s, LEFT, Inches(1.45), W, Inches(1.3))
    _run(ttf.paragraphs[0], title, 34, CHAMPAGNE, bold=True)
    bf = _tb(s, LEFT, Inches(2.9), W, Inches(3.8))
    for i, (lead, rest) in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.space_after = Pt(14)
        _run(p, "— ", 18, accent, bold=True)
        if lead:
            _run(p, lead, 18, CHAMPAGNE, bold=True)
            _run(p, "  " + rest, 18, TEXT)
        else:
            _run(p, rest, 18, TEXT)
    _footer(s, n)
    return s


# ── 1. Title ────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_bg(s)
_mark(s)
tf = _tb(s, LEFT, Inches(2.3), W, Inches(1.6))
_run(tf.paragraphs[0], "CONTXT", 76, CHAMPAGNE, bold=True)
tf2 = _tb(s, LEFT, Inches(3.7), W, Inches(1.0))
_run(tf2.paragraphs[0], "One private context layer for every AI.", 30, GOLD)
tf3 = _tb(s, LEFT, Inches(4.7), W, Inches(1.6))
for i, line in enumerate([
    "It remembers you and acts for you across Claude, ChatGPT & Gemini —",
    "and your private data never leaves your device.",
]):
    p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
    _run(p, line, 18, MUTED)
tf4 = _tb(s, LEFT, Inches(6.4), W, Inches(0.8))
p = tf4.paragraphs[0]
_run(p, "AMD Developer Hackathon ACT II · Track 3 (Unicorn)", 13, MUTED)
p2 = tf4.add_paragraph()
_run(p2, "blackx16.github.io/contxt  ·  github.com/Blackx16/contxt", 13, PATINA)

# ── 2–11 content ──────────────────────────────────────────────────────────────
content(2, "The problem", "Every AI grew a memory. None of them share it.", [
    ("Five walled gardens.", "ChatGPT, Claude, Gemini, Copilot, Grok — each remembers you separately."),
    ("You repeat yourself.", "Re-introducing who you are, what you're building, every time you switch."),
    ("The memory startups", "(Mem0, Supermemory, Letta) are cloud-first and built for developers, not you."),
])
content(3, "The solution", "Your context, everywhere — and it stays yours.", [
    ("Portable.", "One context layer that works across every AI, not locked inside one."),
    ("Private by default.", "A Crown-Jewels Gateway keeps sensitive context on your device."),
    ("Consumer-first.", "The piece nobody built: memory that follows you and protects you."),
], accent=PATINA)
content(4, "How it works", "Two tiers, decided on-device.", [
    ("Crown-Jewels Gateway.", "On-device Gemma 3 270M + deterministic rules classify every item."),
    ("PRIVATE (gold).", "Kept on-device, end-to-end encrypted. The cloud is a blind relay of ciphertext."),
    ("SHARED (patina).", "Distilled into reusable context cards any AI can read over MCP."),
])
content(5, "The product", "A browser extension that plugs into your life.", [
    ("Connect.", "Gmail, Calendar & Notion via in-extension OAuth."),
    ("Tier on-device.", "Every pulled item is classified locally — SHARED shown, PRIVATE kept & counted."),
    ("Your choice.", "On-device mode (WebGPU Gemma) or online-only — nothing private stored."),
])
content(6, "The payoff", "It injects your context into any AI, live.", [
    ("Auto-inject.", "Open Claude / ChatGPT / Gemini — your SHARED context lands in the composer."),
    ("Visible trust.", "A badge shows 'N shared → this AI · P private kept on-device'."),
    ("The guarantee.", "The AI answers with your context and never sees the crown jewels."),
])
content(7, "Privacy, proven", "The cloud never holds your keys.", [
    ("Ciphertext only.", "PRIVATE cards sit in the cloud as opaque AES-256-GCM blobs."),
    ("Key by QR.", "Moving devices? The key crosses device-to-device by QR, never the cloud."),
    ("Decrypt locally.", "A second device pulls the same blob and decrypts it — byte-for-byte."),
])
content(8, "One product", "A web dashboard that mirrors the extension.", [
    ("Deployed & live.", "blackx16.github.io/contxt — a static SvelteKit app on GitHub Pages."),
    ("Demo / Live toggle.", "Explore a built-in demo, or read your real context from the extension."),
    ("Real-time bridge.", "Connect a source in the extension → the site updates instantly."),
], accent=PATINA)
content(9, "Under the hood", "Small, portable, open.", [
    ("On-device:", "Gemma 3 270M via Transformers.js + WebGPU; deterministic-rules safety floor."),
    ("Cloud distill:", "gpt-oss-120B on Fireworks AI (SHARED tier → context cards + draft_reply)."),
    ("Stack:", "SvelteKit · MV3 extension · Python MCP server + HTTP bridge · Web Crypto · public Docker image (GHCR)."),
])
content(10, "On-device model", "Fine-tuning Gemma for the privacy gateway.", [
    ("Tiny + local.", "Gemma 3 270M runs entirely in the browser — no data leaves for tiering."),
    ("Safety floor.", "Deterministic rules catch crown jewels even if the model is unsure."),
    ("In progress.", "A fine-tuned Gemma 270M to sharpen PRIVATE-vs-SHARED decisions on-device."),
])
content(11, "Roadmap", "One context layer. Every AI. Your data stays yours.", [
    ("Today.", "Web + extension, live injection across Claude / ChatGPT / Gemini."),
    ("Next.", "Fully-local mobile (Gemma 3n on-device), Signal-grade multi-device sync."),
    ("Try it.", "blackx16.github.io/contxt · github.com/Blackx16/contxt · ghcr.io/blackx16/contxt"),
])

out = Path(__file__).parent / "contxt-pitch.pptx"
prs.save(out)
print("wrote", out, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
